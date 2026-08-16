import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_detailed_presentation():
    prs = Presentation()
    
    # ── Configurations de Design ──
    def apply_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(18, 18, 30) # Dark violet/blue
        
    def format_title(title_shape, text, size=40):
        if title_shape is None: return
        title_shape.text = text
        for p in title_shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            p.font.bold = True
            p.font.size = Pt(size)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
    def format_content(content_shape, points, size=24):
        if content_shape is None: return
        content_shape.text = "" # Clear default
        tf = content_shape.text_frame
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(size)
            p.font.color.rgb = RGBColor(230, 230, 250)
            p.space_after = Pt(14)

    # ── 1. Page de Garde ──
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_dark_background(slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Soutenance de Projet de Fin d'Études"
    subtitle.text = "DevMob : Plateforme Mobile de Gestion d'Événements\n\n\nPrésenté par : [Votre Nom]"
    
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(180, 180, 255)
        p.font.size = Pt(28)

    # ── 2. Plan ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "Plan de la Présentation")
    points = [
        "1. Contexte et Problématique",
        "2. Solution Proposée (DevMob)",
        "3. Technologies et Architecture",
        "4. Fonctionnalités par Espace",
        "5. Design et Expérience Utilisateur",
        "6. Démonstration",
        "7. Conclusion et Perspectives"
    ]
    format_content(slide.placeholders[1], points, 26)

    # ── 3. Problématique ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "1. Contexte et Problématique")
    points = [
        "Difficulté de trouver et de centraliser les événements locaux.",
        "Processus de réservation de billets souvent complexe et lent.",
        "Manque d'outils analytiques simples pour les organisateurs.",
        "Besoin d'une communication fluide entre organisateurs et participants."
    ]
    format_content(slide.placeholders[1], points, 24)

    # ── 4. Solution ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "2. La Solution : DevMob")
    points = [
        "Une application mobile tout-en-un pour la gestion événementielle.",
        "Trois espaces distincts : Client, Organisateur, et Administrateur.",
        "Génération de billets électroniques (QR Code) pour un accès facile.",
        "Une interface premium, intuitive et très rapide."
    ]
    format_content(slide.placeholders[1], points, 24)

    # ── 5. Technologies ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "3. Technologies Utilisées")
    points = [
        "Frontend Mobile : Flutter (Dart) pour une application multiplateforme fluide.",
        "Backend & Base de données : Firebase (Cloud Firestore, Authentication).",
        "Stockage des médias : Firebase Storage (Images d'événements et de profils).",
        "Cartographie : Flutter Map & OpenStreetMap pour la localisation des événements.",
        "Génération de Code : QR Flutter pour les billets virtuels."
    ]
    format_content(slide.placeholders[1], points, 22)

    # ── 6. Espace Utilisateur ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "4. Espace Utilisateur (Client)")
    points = [
        "Découverte : Recherche par catégorie, filtre et carte interactive.",
        "Réservation : Système de réservation de billets instantané.",
        "Mes Billets : Portefeuille de billets électroniques (QR Code).",
        "Notifications : Alertes en temps réel pour les nouveautés et réservations.",
        "Avis : Possibilité de noter et de laisser des commentaires sur les événements."
    ]
    format_content(slide.placeholders[1], points, 22)

    # ── 7. Espace Organisateur ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "4. Espace Organisateur")
    points = [
        "Création : Formulaire complet pour ajouter de nouveaux événements.",
        "Tableau de Bord : Suivi des événements actifs et des réservations.",
        "Statistiques : Visualisation des revenus et du nombre de billets vendus.",
        "Gestion : Validation des tickets via scan QR (Simulation) et suivi des clients."
    ]
    format_content(slide.placeholders[1], points, 24)

    # ── 8. Espace Admin ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "4. Espace Administrateur")
    points = [
        "Contrôle Global : Vue d'ensemble sur toute la plateforme.",
        "Gestion des Utilisateurs : Possibilité de bloquer ou supprimer des comptes.",
        "Gestion des Événements : Modération du contenu publié par les organisateurs.",
        "Tableau de Bord Global : Suivi de l'évolution de l'application."
    ]
    format_content(slide.placeholders[1], points, 24)

    # ── 9. Design ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "5. Design et Expérience Utilisateur")
    points = [
        "Design System : 'Premium Glassmorphism' avec des effets de flou et de profondeur.",
        "Harmonie Visuelle : Gradients dynamiques, couleurs sombres et élégantes.",
        "Animations : Micro-animations fluides (Flutter Animate) pour une navigation vivante.",
        "Ergonomie : Boutons clairs, icônes arrondies (Material Rounded), interface non surchargée."
    ]
    format_content(slide.placeholders[1], points, 22)

    # ── 10. Demo ──
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_dark_background(slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "6. Démonstration de l'Application"
    subtitle.text = "(Présentation en direct sur l'émulateur / téléphone)"
    title.text_frame.paragraphs[0].font.size = Pt(50)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(160, 160, 255)

    # ── 11. Conclusion ──
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_dark_background(slide)
    format_title(slide.shapes.title, "7. Conclusion et Perspectives")
    points = [
        "Objectifs atteints : Application fonctionnelle, fluide et prête à l'usage.",
        "Perspectives d'avenir :",
        "   - Intégration d'une vraie passerelle de paiement bancaire.",
        "   - Système de recommandation basé sur l'Intelligence Artificielle.",
        "   - Version Web pour l'espace d'administration."
    ]
    format_content(slide.placeholders[1], points, 24)

    # ── 12. Fin ──
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_dark_background(slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Merci pour votre attention !"
    subtitle.text = "DevMob : Avez-vous des questions ?"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(160, 160, 255)

    # ── Sauvegarde ──
    prs.save('DevMob_Presentation_Detaillee.pptx')
    print("Présentation détaillée 'DevMob_Presentation_Detaillee.pptx' créée avec succès !")

if __name__ == '__main__':
    create_detailed_presentation()
